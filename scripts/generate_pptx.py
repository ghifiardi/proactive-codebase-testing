#!/usr/bin/env python3
"""
Generate PowerPoint presentation from documentation.

Requires: python-pptx
Install: pip install python-pptx
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Add project root to path
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

def create_presentation():
    """Create PowerPoint presentation from documentation."""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Proactive Codebase Testing Platform"
    subtitle.text = "AI-Powered Code Analysis\nVersion 0.1.0 | November 2025"
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Executive Summary"
    tf = content.text_frame
    tf.text = "AI-powered code analyzer using Claude API"
    p = tf.add_paragraph()
    p.text = "• Detects security vulnerabilities, bugs, and quality issues"
    p = tf.add_paragraph()
    p.text = "• Supports 20+ programming languages"
    p = tf.add_paragraph()
    p.text = "• Multiple interfaces: CLI, REST API, GitHub Actions"
    p = tf.add_paragraph()
    p.text = "• Production-ready with rate limiting, monitoring, logging"
    
    # Slide 3: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Features"
    tf = content.text_frame
    tf.text = "✅ AI-Powered Analysis"
    p = tf.add_paragraph()
    p.text = "✅ Multi-Language Support (20+ languages)"
    p = tf.add_paragraph()
    p.text = "✅ Multiple Interfaces (CLI, API, CI/CD)"
    p = tf.add_paragraph()
    p.text = "✅ Flexible Output (JSON, HTML, SARIF)"
    p = tf.add_paragraph()
    p.text = "✅ Production Ready (Rate limiting, monitoring)"
    p = tf.add_paragraph()
    p.text = "✅ Comprehensive Testing (>80% coverage)"
    
    # Slide 4: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Architecture"
    tf = content.text_frame
    tf.text = "User Interfaces"
    p = tf.add_paragraph()
    p.text = "  • CLI Tool"
    p = tf.add_paragraph()
    p.text = "  • REST API"
    p = tf.add_paragraph()
    p.text = "  • GitHub Actions"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Core Engine"
    p = tf.add_paragraph()
    p.text = "  • Parser (Language detection)"
    p = tf.add_paragraph()
    p.text = "  • Analyzer (Claude API)"
    p = tf.add_paragraph()
    p.text = "  • Reporters (JSON, HTML, SARIF)"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Production Features"
    p = tf.add_paragraph()
    p.text = "  • Rate Limiting"
    p = tf.add_paragraph()
    p.text = "  • Monitoring & Metrics"
    p = tf.add_paragraph()
    p.text = "  • Security Headers"
    p = tf.add_paragraph()
    p.text = "  • Logging"
    
    # Slide 5: What It Analyzes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "What It Analyzes"
    tf = content.text_frame
    tf.text = "🔒 Security Vulnerabilities"
    p = tf.add_paragraph()
    p.text = "  • SQL Injection, XSS, Authentication issues"
    p = tf.add_paragraph()
    p.text = "  • Hardcoded secrets, Insecure deserialization"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "🐛 Bugs"
    p = tf.add_paragraph()
    p.text = "  • Null pointers, Resource leaks"
    p = tf.add_paragraph()
    p.text = "  • Race conditions, Logic errors"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "✨ Code Quality"
    p = tf.add_paragraph()
    p.text = "  • Code smells, Anti-patterns"
    p = tf.add_paragraph()
    p.text = "  • Best practice violations"
    
    # Slide 6: Usage Examples
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Usage Examples"
    tf = content.text_frame
    tf.text = "CLI:"
    p = tf.add_paragraph()
    p.text = "  python -m src.cli.main analyze . --format html"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "REST API:"
    p = tf.add_paragraph()
    p.text = "  POST /api/analyze"
    p = tf.add_paragraph()
    p.text = "  { \"code\": \"...\", \"language\": \"python\" }"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "GitHub Actions:"
    p = tf.add_paragraph()
    p.text = "  Automatic scanning on push/PR"
    
    # Slide 7: Implementation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Implementation (Langkah 1-12)"
    tf = content.text_frame
    tf.text = "✅ Langkah 1-3: Core Foundation"
    p = tf.add_paragraph()
    p.text = "✅ Langkah 4: Reporters (JSON, HTML, SARIF)"
    p = tf.add_paragraph()
    p.text = "✅ Langkah 5: CLI Interface"
    p = tf.add_paragraph()
    p.text = "✅ Langkah 6: REST API"
    p = tf.add_paragraph()
    p.text = "✅ Langkah 7: GitHub Integration"
    p = tf.add_paragraph()
    p.text = "✅ Langkah 8: Testing (>80% coverage)"
    p = tf.add_paragraph()
    p.text = "✅ Langkah 9: Docker Optimization"
    p = tf.add_paragraph()
    p.text = "✅ Langkah 10: Advanced CI/CD"
    p = tf.add_paragraph()
    p.text = "✅ Langkah 11: Extended Documentation"
    p = tf.add_paragraph()
    p.text = "✅ Langkah 12: Production Hardening"
    
    # Slide 8: Production Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Production Features"
    tf = content.text_frame
    tf.text = "🔒 Security"
    p = tf.add_paragraph()
    p.text = "  • Rate limiting (per-IP)"
    p = tf.add_paragraph()
    p.text = "  • Security headers (XSS, CSP, etc.)"
    p = tf.add_paragraph()
    p.text = "  • Error handling"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "📊 Monitoring"
    p = tf.add_paragraph()
    p.text = "  • Request/response timing"
    p = tf.add_paragraph()
    p.text = "  • Analysis statistics"
    p = tf.add_paragraph()
    p.text = "  • Performance metrics"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "📝 Logging"
    p = tf.add_paragraph()
    p.text = "  • Rotating file handlers"
    p = tf.add_paragraph()
    p.text = "  • Structured logging"
    
    # Slide 9: Deployment Options
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Deployment Options"
    tf = content.text_frame
    tf.text = "🐳 Docker"
    p = tf.add_paragraph()
    p.text = "  docker-compose up"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "☁️ Cloud Platforms"
    p = tf.add_paragraph()
    p.text = "  • Heroku, AWS, GCP, Azure"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "☸️ Kubernetes"
    p = tf.add_paragraph()
    p.text = "  • Full K8s manifests included"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "💻 Local"
    p = tf.add_paragraph()
    p.text = "  python -m src.cli.main analyze ."
    
    # Slide 10: Statistics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Statistics"
    tf = content.text_frame
    tf.text = "📊 Code"
    p = tf.add_paragraph()
    p.text = "  • 21 Python files"
    p = tf.add_paragraph()
    p.text = "  • ~4,600 lines of code"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "✅ Testing"
    p = tf.add_paragraph()
    p.text = "  • >80% code coverage"
    p = tf.add_paragraph()
    p.text = "  • 40+ unit tests"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "📚 Documentation"
    p = tf.add_paragraph()
    p.text = "  • 3,500+ lines of docs"
    p = tf.add_paragraph()
    p.text = "  • Complete API reference"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "🚀 CI/CD"
    p = tf.add_paragraph()
    p.text = "  • 5 GitHub Actions workflows"
    
    # Slide 11: Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Next Steps"
    tf = content.text_frame
    tf.text = "1. Get Anthropic API Key"
    p = tf.add_paragraph()
    p.text = "2. Install dependencies"
    p = tf.add_paragraph()
    p.text = "3. Run first analysis"
    p = tf.add_paragraph()
    p.text = "4. Integrate with CI/CD"
    p = tf.add_paragraph()
    p.text = "5. Deploy to production"
    
    # Slide 12: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You"
    subtitle.text = "Proactive Codebase Testing Platform\n\nMade with ❤️ for secure, clean code\n\nDocumentation: docs/\nExamples: examples/\n\nVersion 0.1.0 | November 2025"
    
    # Save presentation
    output_path = project_root / "documentation_output" / "Proactive_Codebase_Testing_Platform.pptx"
    output_path.parent.mkdir(exist_ok=True)
    prs.save(str(output_path))
    
    print(f"✅ PowerPoint presentation created: {output_path}")
    return output_path

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ python-pptx not installed.")
        print("   Install: pip install python-pptx")
        sys.exit(1)
    except Exception as e:
        print(f"❌ Error: {e}")
        sys.exit(1)

